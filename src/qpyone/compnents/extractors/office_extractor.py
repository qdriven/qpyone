import docx2txt
import csv
import pptx

from qpyone.compnents.extractors.models import TextContextExtractor, validate_mimetype


class WordTextExtractor(TextContextExtractor):

    def extract_text_from_filepath(self, file_path: str) -> str:
        self.validate_file_type(file_path)
        with open(file_path, "rb") as f:
            extracted_text = docx2txt.process(f)
        return extracted_text

    def validate_file_type(self, file_path: str):
        # Get the mimetype of the file based on its extension
        validate_mimetype(file_path, [".doc", ".docx"],
                          "application/vnd.openxmlformats-officedocument.wordprocessingml.document")


class PPTTextExtractor(TextContextExtractor):

    def extract_text_from_filepath(self, file_path: str) -> str:
        self.validate_file_type(file_path)
        extracted_text = ""
        with open(file_path, "rb") as f:
            presentation = pptx.Presentation(f)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                extracted_text += run.text + " "
                        extracted_text += "\n"
        return extracted_text

    def validate_file_type(self, file_path: str):
        # Get the mimetype of the file based on its extension
        validate_mimetype(file_path, [".ppt", ".pptx"],
                          "application/vnd.openxmlformats-officedocument.presentationml.presentation")


class CSVTextExtractor(TextContextExtractor):

    def extract_text_from_filepath(self, file_path: str) -> str:
        self.validate_file_type(file_path)
        extracted_text = ""
        with open(file_path, "rb") as f:
            decoded_buffer = (line.decode("utf-8") for line in f)
            reader = csv.reader(decoded_buffer)
            for row in reader:
                extracted_text += " ".join(row) + "\n"
        return extracted_text

    def validate_file_type(self, file_path: str):
        # Get the mimetype of the file based on its extension
        validate_mimetype(file_path, [".csv", ".xls", ".xlsx"],
                          "text/csv")
